#!/usr/bin/env python3
"""
创建5个风格迥异的测试PPT，用于完整性阶段测试
覆盖：商业/科技/教育/渐变/极简等不同风格
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import nsmap
from lxml import etree
import os

# 确保命名空间正确
PRS_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

def hex2rgb(hex_s):
    """将6位HEX转为RGBColor(r,g,b)"""
    h = hex_s.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_animation(slide, shape, anim_type="appear"):
    """为形状添加动画"""
    sp_tree = slide.shapes._spTree
    
    timing_nodes = sp_tree.findall(f".//{{{PRS_NS}}}timing")
    if not timing_nodes:
        timing = etree.SubElement(sp_tree, f"{{{PRS_NS}}}timing")
        tnLst = etree.SubElement(timing, f"{{{PRS_NS}}}tnLst")
        par = etree.SubElement(tnLst, f"{{{PRS_NS}}}par")
        cTn = etree.SubElement(par, f"{{{PRS_NS}}}cTn")
        cTn.set("id", "1")
        cTn.set("dur", "indefinite")
        cTn.set("restart", "never")
        childTnLst = etree.SubElement(cTn, f"{{{PRS_NS}}}childTnLst")
        seq = etree.SubElement(childTnLst, f"{{{PRS_NS}}}seq")
        seq.set("concurrent", "1")
        seq.set("nextAc", "seek")
        cTn2 = etree.SubElement(seq, f"{{{PRS_NS}}}cTn")
        cTn2.set("id", "2")
        cTn2.set("dur", "indefinite")
        cTn2.set("nodeType", "mainSeq")
    
    timing = sp_tree.find(f"{{{PRS_NS}}}timing")
    tnLst = timing.find(f"{{{PRS_NS}}}tnLst")
    par = tnLst.find(f"{{{PRS_NS}}}par")
    cTn = par.find(f"{{{PRS_NS}}}cTn")
    childTnLst = cTn.find(f"{{{PRS_NS}}}childTnLst")
    seq = childTnLst.find(f"{{{PRS_NS}}}seq")
    cTn2 = seq.find(f"{{{PRS_NS}}}cTn")
    childTnLst2 = cTn2.find(f"{{{PRS_NS}}}childTnLst")
    if childTnLst2 is None:
        childTnLst2 = etree.SubElement(cTn2, f"{{{PRS_NS}}}childTnLst")
    
    par_anim = etree.SubElement(childTnLst2, f"{{{PRS_NS}}}par")
    cTn_anim = etree.SubElement(par_anim, f"{{{PRS_NS}}}cTn")
    cTn_anim.set("id", str(100 + len(childTnLst2)))
    cTn_anim.set("fill", "hold")
    stCondLst = etree.SubElement(cTn_anim, f"{{{PRS_NS}}}stCondLst")
    cond = etree.SubElement(stCondLst, f"{{{PRS_NS}}}cond")
    cond.set("delay", "indefinite")
    cond.set("evt", "onClick")
    childTnLst3 = etree.SubElement(cTn_anim, f"{{{PRS_NS}}}childTnLst")
    par2 = etree.SubElement(childTnLst3, f"{{{PRS_NS}}}par")
    cTn3 = etree.SubElement(par2, f"{{{PRS_NS}}}cTn")
    cTn3.set("id", str(200 + len(childTnLst2)))
    cTn3.set("fill", "hold")
    cTn3.set("presetClass", "entr")
    cTn3.set("presetID", "1")
    cTn3.set("presetSubtype", "0")
    stCondLst2 = etree.SubElement(cTn3, f"{{{PRS_NS}}}stCondLst")
    cond2 = etree.SubElement(stCondLst2, f"{{{PRS_NS}}}cond")
    cond2.set("delay", "0")
    childTnLst4 = etree.SubElement(cTn3, f"{{{PRS_NS}}}childTnLst")
    set_anim = etree.SubElement(childTnLst4, f"{{{PRS_NS}}}set")
    cBhvr = etree.SubElement(set_anim, f"{{{PRS_NS}}}cBhvr")
    cBhvr.set("override", "childStyle")
    cTn4 = etree.SubElement(cBhvr, f"{{{PRS_NS}}}cTn")
    cTn4.set("id", str(300 + len(childTnLst2)))
    cTn4.set("dur", "500")
    cTn4.set("fill", "hold")
    tgtEl = etree.SubElement(cBhvr, f"{{{PRS_NS}}}tgtEl")
    spTgt = etree.SubElement(tgtEl, f"{{{PRS_NS}}}spTgt")
    spTgt.set("spid", str(shape.shape_id))
    attrNameLst = etree.SubElement(cBhvr, f"{{{PRS_NS}}}attrNameLst")
    attrName = etree.SubElement(attrNameLst, f"{{{PRS_NS}}}attrName")
    attrName.text = "style.visibility"
    to = etree.SubElement(set_anim, f"{{{PRS_NS}}}to")
    strVal = etree.SubElement(to, f"{{{PRS_NS}}}strVal")
    strVal.set("val", "visible")


def add_transition(slide, trans_type="fade"):
    """为幻灯片添加转场效果"""
    transition = etree.SubElement(slide._element, f"{{{PRS_NS}}}transition")
    transition.set("spd", "med")
    
    if trans_type == "fade":
        etree.SubElement(transition, f"{{{PRS_NS}}}fade")
    elif trans_type == "push":
        etree.SubElement(transition, f"{{{PRS_NS}}}push")
    elif trans_type == "wipe":
        etree.SubElement(transition, f"{{{PRS_NS}}}wipe")
    elif trans_type == "cover":
        etree.SubElement(transition, f"{{{PRS_NS}}}cover")


# ═══════════════════════════════════════════════════════════════
# 测试PPT 1：商务深蓝（经典商务风格，6页）
# ═══════════════════════════════════════════════════════════════
def create_business_blue():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    C = {
        "primary": "1B3A5C", "secondary": "5B8DB8", "accent": "E8913A",
        "bg_dark": "1B3A5C", "bg_light": "F8F9FA", "text_dark": "2D3748", "text_light": "FFFFFF",
        "deco": "B0C4DE"
    }
    
    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex2rgb(C["bg_dark"])
    
    title = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1))
    p = title.text_frame.paragraphs[0]
    p.text = "2026年度战略规划"; p.font.size = Pt(44); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["text_light"]); p.alignment = PP_ALIGN.CENTER
    
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.6))
    p = subtitle.text_frame.paragraphs[0]
    p.text = "战略部 · 2026年Q1"; p.font.size = Pt(20)
    p.font.color.rgb = hex2rgb(C["deco"]); p.alignment = PP_ALIGN.CENTER
    
    line = slide.shapes.add_shape(1, Inches(3.5), Inches(3.1), Inches(3), Inches(0.02))
    line.fill.solid(); line.fill.fore_color.rgb = hex2rgb(C["accent"])
    line.line.color.rgb = hex2rgb(C["accent"])
    add_animation(slide, title); add_transition(slide, "fade")
    
    # 目录
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_light"])
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = "目录 CONTENTS"; p.font.size = Pt(32); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["primary"])
    
    items = ["市场环境分析", "战略方向", "执行计划", "风险管控"]
    for i, item in enumerate(items):
        y = 1.5 + i * 0.9
        circle = slide.shapes.add_shape(9, Inches(1.2), Inches(y), Inches(0.5), Inches(0.5))
        circle.fill.solid(); circle.fill.fore_color.rgb = hex2rgb(C["primary"])
        num = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(0.5), Inches(0.5))
        p = num.text_frame.paragraphs[0]
        p.text = str(i+1); p.font.size = Pt(16); p.font.bold = True
        p.font.color.rgb = hex2rgb(C["text_light"]); p.alignment = PP_ALIGN.CENTER
        tb = slide.shapes.add_textbox(Inches(2.0), Inches(y), Inches(6), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = item; p.font.size = Pt(20); p.font.color.rgb = hex2rgb(C["text_dark"])
    add_transition(slide, "push")
    
    # 章节页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_dark"])
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = "市场环境分析"; p.font.size = Pt(40); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["text_light"]); p.alignment = PP_ALIGN.CENTER
    
    line = slide.shapes.add_shape(1, Inches(3.5), Inches(1.8), Inches(3), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = hex2rgb(C["accent"])
    line.line.color.rgb = hex2rgb(C["accent"])
    add_transition(slide, "fade")
    
    # 内容页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_light"])
    
    header = slide.shapes.add_shape(1, 0, 0, Inches(10), Inches(1.0))
    header.fill.solid(); header.fill.fore_color.rgb = hex2rgb(C["primary"])
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = "市场规模与趋势"; p.font.size = Pt(28); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["text_light"])
    
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3.5))
    body.text_frame.word_wrap = True
    pts = ["• 全球市场规模预计达 $2.8万亿，年增长率 12%",
           "• 亚太地区贡献 45% 的增长份额",
           "• 数字化转型加速，企业级需求持续上升",
           "• 竞争格局趋于稳定，头部效应明显"]
    for i, pt in enumerate(pts):
        p = body.text_frame.paragraphs[0] if i == 0 else body.text_frame.add_paragraph()
        p.text = pt; p.font.size = Pt(16); p.font.color.rgb = hex2rgb(C["text_dark"])
        p.space_after = Pt(12)
    add_animation(slide, body); add_transition(slide, "wipe")
    
    # 数据页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_light"])
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = "关键数据指标"; p.font.size = Pt(28); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["primary"])
    
    metrics = [("$2.8T", "市场规模", "同比增长12%", C["primary"]),
               ("45%", "亚太占比", "全球最高", C["secondary"]),
               ("99.9%", "系统可用", "全年稳定", C["accent"])]
    for i, (val, label, desc, color) in enumerate(metrics):
        x = 0.8 + i * 3.0
        card = slide.shapes.add_shape(1, Inches(x), Inches(1.8), Inches(2.5), Inches(2.5))
        card.fill.solid(); card.fill.fore_color.rgb = hex2rgb("FFFFFF")
        val_box = slide.shapes.add_textbox(Inches(x), Inches(2.0), Inches(2.5), Inches(0.8))
        p = val_box.text_frame.paragraphs[0]
        p.text = val; p.font.size = Pt(32); p.font.bold = True
        p.font.color.rgb = hex2rgb(color); p.alignment = PP_ALIGN.CENTER
        label_box = slide.shapes.add_textbox(Inches(x), Inches(2.8), Inches(2.5), Inches(0.5))
        p = label_box.text_frame.paragraphs[0]
        p.text = label; p.font.size = Pt(16); p.font.color.rgb = hex2rgb(C["text_dark"])
        p.alignment = PP_ALIGN.CENTER
    add_transition(slide, "fade")
    
    # 结尾页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_dark"])
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = "感谢聆听"; p.font.size = Pt(40); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["text_light"]); p.alignment = PP_ALIGN.CENTER
    
    contact = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.5))
    p = contact.text_frame.paragraphs[0]
    p.text = "yzc-cn@qq.com"; p.font.size = Pt(14)
    p.font.color.rgb = hex2rgb(C["deco"]); p.alignment = PP_ALIGN.CENTER
    add_transition(slide, "cover")
    
    prs.save("/data/workspace/skill/skills/pptx-template-learner/test_ppts/test_business_blue.pptx")
    print("✅ 商务深蓝PPT已创建（6页）")


# ═══════════════════════════════════════════════════════════════
# 测试PPT 2：暗色科技风（4页）
# ═══════════════════════════════════════════════════════════════
def create_dark_tech():
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
    
    C = {"primary": "00D4FF", "secondary": "1A1A3E", "accent": "FF6B35",
         "bg_dark": "0A0A1A", "bg_light": "1A1A3E", "text_dark": "2D2D5A", "text_light": "E0E0FF"}
    
    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_dark"])
    
    deco = slide.shapes.add_shape(1, Inches(0), Inches(4.5), Inches(10), Inches(1.125))
    deco.fill.solid(); deco.fill.fore_color.rgb = hex2rgb(C["primary"])
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    p = title.text_frame.paragraphs[0]
    p.text = "AI TECH SUMMIT 2026"; p.font.size = Pt(48); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["primary"]); p.alignment = PP_ALIGN.CENTER
    
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.6))
    p = subtitle.text_frame.paragraphs[0]
    p.text = "探索人工智能的无限可能"; p.font.size = Pt(22)
    p.font.color.rgb = hex2rgb(C["text_light"]); p.alignment = PP_ALIGN.CENTER
    add_animation(slide, title); add_transition(slide, "fade")
    
    # 亮色章节页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["primary"])
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "大模型技术演进"; p.font.size = Pt(42); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["bg_dark"]); p.alignment = PP_ALIGN.CENTER
    add_transition(slide, "push")
    
    # 内容页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_light"])
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = "模型规模对比"; p.font.size = Pt(30); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["primary"])
    
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3.5))
    body.text_frame.word_wrap = True
    pts = ["• GPT-5: 1.8T 参数，多模态原生架构",
           "• Claude Opus 4.7: 2.1T 参数，推理能力突破",
           "• Gemini Ultra 2: 1.5T 参数，长上下文128K",
           "• DeepSeek-V4: 671B 参数，成本效率领先"]
    for i, pt in enumerate(pts):
        p = body.text_frame.paragraphs[0] if i == 0 else body.text_frame.add_paragraph()
        p.text = pt; p.font.size = Pt(16); p.font.color.rgb = hex2rgb(C["text_light"])
        p.space_after = Pt(14)
    add_animation(slide, body); add_transition(slide, "wipe")
    
    # 结尾页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_dark"])
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "THANK YOU"; p.font.size = Pt(44); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["primary"]); p.alignment = PP_ALIGN.CENTER
    add_transition(slide, "fade")
    
    prs.save("/data/workspace/skill/skills/pptx-template-learner/test_ppts/test_dark_tech.pptx")
    print("✅ 暗色科技风PPT已创建（4页）")


# ═══════════════════════════════════════════════════════════════
# 测试PPT 3：教育培训绿（5页）
# ═══════════════════════════════════════════════════════════════
def create_education_green():
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
    
    C = {"primary": "2D7A46", "secondary": "8FBC8F", "accent": "E8B923",
         "bg_dark": "2D7A46", "bg_light": "F0FDF4", "text_dark": "1A1A1A", "text_light": "FFFFFF"}
    
    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_dark"])
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1))
    p = title.text_frame.paragraphs[0]
    p.text = "Python编程入门"; p.font.size = Pt(44); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["text_light"]); p.alignment = PP_ALIGN.CENTER
    
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9), Inches(0.6))
    p = subtitle.text_frame.paragraphs[0]
    p.text = "零基础到实战"; p.font.size = Pt(22)
    p.font.color.rgb = hex2rgb(C["secondary"]); p.alignment = PP_ALIGN.CENTER
    add_animation(slide, title); add_transition(slide, "fade")
    
    # 目录
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_light"])
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = "课程大纲"; p.font.size = Pt(32); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["primary"])
    
    items = ["基础语法", "数据结构", "函数与模块", "实战项目"]
    for i, item in enumerate(items):
        y = 1.5 + i * 0.9
        num_box = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(0.5), Inches(0.5))
        p = num_box.text_frame.paragraphs[0]
        p.text = f"0{i+1}"; p.font.size = Pt(18); p.font.bold = True
        p.font.color.rgb = hex2rgb(C["accent"])
        tb = slide.shapes.add_textbox(Inches(2.2), Inches(y), Inches(6), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = item; p.font.size = Pt(20); p.font.color.rgb = hex2rgb(C["text_dark"])
    add_transition(slide, "push")
    
    # 内容页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_light"])
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = "基础语法要点"; p.font.size = Pt(28); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["primary"])
    
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3.5))
    body.text_frame.word_wrap = True
    pts = ["• 变量与数据类型：int, float, str, bool, list, dict",
           "• 控制流：if/else, for, while, 列表推导式",
           "• 函数定义：def, lambda, 参数传递方式",
           "• 异常处理：try/except, raise, finally"]
    for i, pt in enumerate(pts):
        p = body.text_frame.paragraphs[0] if i == 0 else body.text_frame.add_paragraph()
        p.text = pt; p.font.size = Pt(16); p.font.color.rgb = hex2rgb(C["text_dark"])
        p.space_after = Pt(12)
    add_animation(slide, body); add_transition(slide, "wipe")
    
    # 章节页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_dark"])
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "实战项目"; p.font.size = Pt(40); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["text_light"]); p.alignment = PP_ALIGN.CENTER
    add_transition(slide, "fade")
    
    # 结尾页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_dark"])
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "谢谢大家"; p.font.size = Pt(40); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["text_light"]); p.alignment = PP_ALIGN.CENTER
    add_transition(slide, "cover")
    
    prs.save("/data/workspace/skill/skills/pptx-template-learner/test_ppts/test_education_green.pptx")
    print("✅ 教育培训PPT已创建（5页）")


# ═══════════════════════════════════════════════════════════════
# 测试PPT 4：渐变浪漫风粉紫（3页）
# ═══════════════════════════════════════════════════════════════
def create_gradient_romantic():
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
    
    C = {"primary": "9333EA", "secondary": "F472B6", "accent": "FBBF24",
         "bg_dark": "4C1D95", "bg_light": "FDF2F8", "text_dark": "4B0082", "text_light": "FFFFFF"}
    
    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_dark"])
    
    deco = slide.shapes.add_shape(1, Inches(6), Inches(0), Inches(4), Inches(5.625))
    deco.fill.solid(); deco.fill.fore_color.rgb = hex2rgb(C["primary"])
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5.5), Inches(1))
    p = title.text_frame.paragraphs[0]
    p.text = "品牌视觉设计"; p.font.size = Pt(42); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["text_light"])
    
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(5.5), Inches(0.6))
    p = subtitle.text_frame.paragraphs[0]
    p.text = "Brand Visual Design 2026"; p.font.size = Pt(18)
    p.font.color.rgb = hex2rgb(C["secondary"])
    add_animation(slide, title); add_transition(slide, "fade")
    
    # 内容页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_light"])
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = "设计理念"; p.font.size = Pt(30); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["primary"])
    
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3.5))
    body.text_frame.word_wrap = True
    pts = ["• 色彩表达：以紫色系为主调，传递优雅与创意",
           "• 字体选择：使用圆润无衬线字体，体现亲和力",
           "• 图形风格：几何图形+柔和渐变，现代感与温度并存",
           "• 排版逻辑：大面积留白，内容聚焦"]
    for i, pt in enumerate(pts):
        p = body.text_frame.paragraphs[0] if i == 0 else body.text_frame.add_paragraph()
        p.text = pt; p.font.size = Pt(16); p.font.color.rgb = hex2rgb(C["text_dark"])
        p.space_after = Pt(14)
    add_animation(slide, body); add_transition(slide, "wipe")
    
    # 结尾页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_dark"])
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "感谢支持"; p.font.size = Pt(40); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["text_light"]); p.alignment = PP_ALIGN.CENTER
    add_transition(slide, "cover")
    
    prs.save("/data/workspace/skill/skills/pptx-template-learner/test_ppts/test_gradient_romantic.pptx")
    print("✅ 渐变浪漫风PPT已创建（3页）")


# ═══════════════════════════════════════════════════════════════
# 测试PPT 5：极简黑白（3页）
# ═══════════════════════════════════════════════════════════════
def create_minimal_bw():
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
    
    C = {"primary": "000000", "secondary": "808080", "accent": "FF0000",
         "bg_dark": "000000", "bg_light": "FFFFFF", "text_dark": "000000", "text_light": "FFFFFF"}
    
    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_dark"])
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1))
    p = title.text_frame.paragraphs[0]
    p.text = "MINIMAL"; p.font.size = Pt(60); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["text_light"]); p.alignment = PP_ALIGN.CENTER
    
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.5))
    p = subtitle.text_frame.paragraphs[0]
    p.text = "Less is more"; p.font.size = Pt(18)
    p.font.color.rgb = hex2rgb(C["secondary"]); p.alignment = PP_ALIGN.CENTER
    add_animation(slide, title); add_transition(slide, "fade")
    
    # 内容页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_light"])
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    p = tb.text_frame.paragraphs[0]
    p.text = "设计原则"; p.font.size = Pt(32); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["text_dark"])
    
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
    body.text_frame.word_wrap = True
    pts = ["• 去繁从简：每个元素必须有存在的理由",
           "• 留白是内容：空白区域传递呼吸感",
           "• 对比产生力量：黑与白的极致对比",
           "• 字体即设计：Helvetica Neue，极致简约"]
    for i, pt in enumerate(pts):
        p = body.text_frame.paragraphs[0] if i == 0 else body.text_frame.add_paragraph()
        p.text = pt; p.font.size = Pt(16); p.font.color.rgb = hex2rgb(C["text_dark"])
        p.space_after = Pt(16)
    
    accent_line = slide.shapes.add_shape(1, Inches(0.5), Inches(4.5), Inches(2), Inches(0.05))
    accent_line.fill.solid(); accent_line.fill.fore_color.rgb = hex2rgb(C["accent"])
    accent_line.line.color.rgb = hex2rgb(C["accent"])
    add_animation(slide, body); add_transition(slide, "push")
    
    # 结尾页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = hex2rgb(C["bg_dark"])
    
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "THE END"; p.font.size = Pt(48); p.font.bold = True
    p.font.color.rgb = hex2rgb(C["text_light"]); p.alignment = PP_ALIGN.CENTER
    add_transition(slide, "fade")
    
    prs.save("/data/workspace/skill/skills/pptx-template-learner/test_ppts/test_minimal_bw.pptx")
    print("✅ 极简黑白PPT已创建（3页）")


if __name__ == "__main__":
    create_business_blue()
    create_dark_tech()
    create_education_green()
    create_gradient_romantic()
    create_minimal_bw()
    print("\n🎉 全部5个测试PPT创建完成！")
